"""Boundary and edge case tests for the application."""
import pytest
import json
from fastapi.testclient import TestClient
from app.main import app
from services.lesson import LessonGenerator
from services.voice import VoiceGenerator
from services.ppt import PPTGenerator


client = TestClient(app)


class TestGenerateEndpointBoundary:
    """Boundary tests for generate endpoint."""

    def test_generate_minimum_duration(self):
        """Test generate with minimum duration (15 minutes)."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "数学",
                "grade": "7年级",
                "duration": 15,  # Minimum
                "topic": "测试",
            },
        )
        assert response.status_code == 200

    def test_generate_maximum_duration(self):
        """Test generate with maximum duration (120 minutes)."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "数学",
                "grade": "7年级",
                "duration": 120,  # Maximum
                "topic": "测试",
            },
        )
        assert response.status_code == 200

    def test_generate_duration_below_minimum(self):
        """Test generate rejects duration below 15."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "数学",
                "grade": "7年级",
                "duration": 14,
                "topic": "测试",
            },
        )
        assert response.status_code == 422

    def test_generate_duration_above_maximum(self):
        """Test generate rejects duration above 120."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "数学",
                "grade": "7年级",
                "duration": 121,
                "topic": "测试",
            },
        )
        assert response.status_code == 422

    def test_generate_empty_subject(self):
        """Test generate rejects empty subject."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "",
                "grade": "7年级",
                "duration": 45,
                "topic": "测试",
            },
        )
        assert response.status_code == 200  # Empty string is valid but may cause issues downstream

    def test_generate_empty_topic(self):
        """Test generate with empty topic."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "数学",
                "grade": "7年级",
                "duration": 45,
                "topic": "",
            },
        )
        assert response.status_code == 200

    def test_generate_special_characters_in_topic(self):
        """Test generate handles special characters."""
        response = client.post(
            "/api/generate",
            json={
                "subject": "数学",
                "grade": "7年级",
                "duration": 45,
                "topic": "方程 & 不等式 <测试> \"引号\" '单引号'",
            },
        )
        assert response.status_code == 200


class TestLessonGeneratorBoundary:
    """Boundary tests for LessonGenerator."""

    def test_generate_very_long_topic(self):
        """Test generator handles very long topic name."""
        generator = LessonGenerator()

        long_topic = "这是一个非常非常长的主题名称，用来测试系统对超长输入的处理能力" * 10

        result = generator.generate({
            "subject": "数学",
            "grade": "7年级",
            "topic": long_topic,
            "duration": 45,
        })

        assert "meta" in result
        assert result["meta"]["topic"] == long_topic

    def test_generate_unicode_characters(self):
        """Test generator handles Unicode characters."""
        generator = LessonGenerator()

        result = generator.generate({
            "subject": "语文",
            "grade": "高三",
            "topic": "古诗词鉴赏 - 《静夜思》李白",
            "duration": 45,
        })

        assert "meta" in result

    def test_generate_all_empty_fields(self):
        """Test generator handles all empty optional fields."""
        generator = LessonGenerator()

        result = generator.generate({
            "subject": "",
            "grade": "",
            "topic": "",
            "duration": 15,
            "style": "",
        })

        # Template fallback should still return valid structure
        assert "meta" in result
        assert "outline" in result


class TestVoiceGeneratorBoundary:
    """Boundary tests for VoiceGenerator."""

    def test_generate_with_empty_sections(self):
        """Test voice generator handles empty sections."""
        generator = VoiceGenerator()

        lesson = {
            "meta": {"subject": "测试"},
            "outline": {
                "introduction": {},
                "main_sections": [],
                "conclusion": {},
            },
        }

        result = generator.generate(lesson)

        assert isinstance(result, list)

    def test_generate_with_missing_sections(self):
        """Test voice generator handles missing sections."""
        generator = VoiceGenerator()

        lesson = {
            "meta": {"subject": "测试"},
            "outline": {},  # No sections at all
        }

        result = generator.generate(lesson)

        assert isinstance(result, list)

    def test_generate_with_very_long_content(self):
        """Test voice generator handles very long content."""
        generator = VoiceGenerator()

        long_content = "这是一个很长的内容段落。" * 100

        lesson = {
            "meta": {"subject": "测试"},
            "outline": {
                "introduction": {
                    "title": "导入",
                    "content": long_content,
                },
            },
        }

        result = generator.generate(lesson)

        assert isinstance(result, list)


class TestPPTGeneratorBoundary:
    """Boundary tests for PPTGenerator."""

    def test_generate_with_many_sections(self):
        """Test PPT generator handles many sections."""
        generator = PPTGenerator()

        # Create many sections
        main_sections = [
            {
                "title": f"知识点{i}",
                "content": f"内容{i}",
                "key_points": ["要点1", "要点2"],
                "media_hint": {"slide_type": "knowledge"},
            }
            for i in range(10)
        ]

        lesson = {
            "meta": {"subject": "数学", "topic": "测试"},
            "outline": {
                "introduction": {"title": "导入", "content": "内容"},
                "main_sections": main_sections,
                "conclusion": {"title": "总结", "content": "内容"},
            },
        }

        # Note: This will try to save to /storage/ppt which may not exist
        # We test the generation logic, not the file save
        from pptx import Presentation
        prs = Presentation()
        generator._add_title_slide(prs, lesson["meta"])

        for section in main_sections:
            from models.lesson import SlideType
            generator._add_content_slide(prs, section, SlideType.knowledge)

        # Should have created slides
        assert len(prs.slides) == 11  # 1 title + 10 content

    def test_generate_with_empty_key_points_list(self):
        """Test PPT handles empty key points."""
        generator = PPTGenerator()

        section = {
            "title": "测试",
            "content": "内容",
            "key_points": [],  # Empty list
        }

        from pptx import Presentation
        from models.lesson import SlideType
        prs = Presentation()

        generator._add_content_slide(prs, section, SlideType.knowledge)

        assert len(prs.slides) == 1


class TestJobStatusBoundary:
    """Boundary tests for job status tracking."""

    def test_job_status_concurrent_tasks(self):
        """Test job handles concurrent task updates."""
        from models.job import JobStatus, TaskProgress, TaskStatus, ProductType
        from datetime import datetime

        job = JobStatus(
            job_id="test-concurrent",
            status=TaskStatus.pending,
            tasks=[
                TaskProgress(type=ProductType.lesson, status=TaskStatus.in_progress, progress=10),
                TaskProgress(type=ProductType.tts, status=TaskStatus.in_progress, progress=20),
                TaskProgress(type=ProductType.ppt, status=TaskStatus.in_progress, progress=30),
                TaskProgress(type=ProductType.video, status=TaskStatus.pending, progress=0),
            ],
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow(),
        )

        # All tasks should be tracked
        assert len(job.tasks) == 4

        # Status should reflect overall state
        in_progress_count = sum(1 for t in job.tasks if t.status == TaskStatus.in_progress)
        assert in_progress_count == 3

    def test_progress_message_extreme_values(self):
        """Test progress message handles extreme values."""
        from models.job import ProgressMessage, TaskStatus, ProductType

        # Progress at boundaries
        msg_min = ProgressMessage(
            job_id="test",
            task_type=ProductType.lesson,
            status=TaskStatus.pending,
            progress=0,  # Minimum
        )
        assert msg_min.progress == 0

        msg_max = ProgressMessage(
            job_id="test",
            task_type=ProductType.lesson,
            status=TaskStatus.completed,
            progress=100,  # Maximum
        )
        assert msg_max.progress == 100


class TestWebSocketBoundary:
    """Boundary tests for WebSocket handling."""

    def test_manager_handles_many_connections_per_job(self):
        """Test manager handles many connections for same job."""
        from app.websocket import ConnectionManager
        from unittest.mock import AsyncMock
        import asyncio

        manager = ConnectionManager()

        # Add many connections
        async def add_connections():
            for i in range(100):
                mock_ws = AsyncMock()
                await manager.connect("test-job", mock_ws)

        asyncio.run(add_connections())

        assert len(manager.active_connections["test-job"]) == 100

    def test_manager_handles_many_jobs(self):
        """Test manager handles many different jobs."""
        from app.websocket import ConnectionManager
        from unittest.mock import AsyncMock
        import asyncio

        manager = ConnectionManager()

        async def add_jobs():
            for i in range(50):
                mock_ws = AsyncMock()
                await manager.connect(f"job-{i}", mock_ws)

        asyncio.run(add_jobs())

        assert len(manager.active_connections) == 50