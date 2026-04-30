"""PPT generation service using python-pptx."""
import os
import uuid
from pathlib import Path
from pptx import Presentation
from config.settings import settings
from models.lesson import SlideType

# Storage directory for PPT files
STORAGE_DIR = Path(os.getenv("STORAGE_DIR", "/tmp/yaeteaching/storage"))
PPT_DIR = STORAGE_DIR / "ppt"


class PPTGenerator:
    """Generate PowerPoint from lesson content."""

    def __init__(self):
        self._ensure_storage_dir()
        self.template_path = settings.minio_bucket  # TODO: actual template storage

    def _ensure_storage_dir(self):
        """Ensure storage directory exists."""
        PPT_DIR.mkdir(parents=True, exist_ok=True)

    def generate(self, lesson_json: dict) -> str:
        """Generate PPT file from lesson JSON.

        Args:
            lesson_json: Lesson JSON dict

        Returns:
            Path to generated PPT file
        """
        prs = Presentation()
        outline = lesson_json.get("outline", {})
        meta = lesson_json.get("meta", {})

        # Title slide
        self._add_title_slide(prs, meta)

        # Introduction
        intro = outline.get("introduction", {})
        if intro:
            self._add_content_slide(prs, intro, SlideType.title)

        # Main sections
        for section in outline.get("main_sections", []):
            slide_type = section.get("media_hint", {}).get("slide_type", "knowledge")
            try:
                st = SlideType(slide_type)
            except ValueError:
                st = SlideType.knowledge
            self._add_content_slide(prs, section, st)

        # Conclusion
        conclusion = outline.get("conclusion", {})
        if conclusion:
            self._add_content_slide(prs, conclusion, SlideType.summary)

        # Generate unique filename
        job_id = str(uuid.uuid4())[:8]
        topic = meta.get("topic", "lesson")
        # Sanitize topic for filename
        safe_topic = "".join(c if c.isalnum() or c in "-_" else "_" for c in topic)

        # Save to storage
        output_path = PPT_DIR / f"{safe_topic}_{job_id}.pptx"
        prs.save(str(output_path))

        return str(output_path)

    def _add_title_slide(self, prs: Presentation, meta: dict):
        """Add title slide with lesson metadata."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = meta.get("topic", "课程")

        subtitle = slide.placeholders[1]
        subtitle.text = f"{meta.get('subject', '')} - {meta.get('grade', '')}"

    def _add_content_slide(self, prs: Presentation, section: dict, slide_type: SlideType):
        """Add content slide based on section and type."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = section.get("title", "")

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = section.get("content", "")

        # Add key points as bullet list
        for point in section.get("key_points", []):
            p = tf.add_paragraph()
            p.text = point
            p.level = 1